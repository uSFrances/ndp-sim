import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from estimate_ttft_from_config import build_estimate


ROOT = Path(__file__).resolve().parents[1]
CONFIG_PATH = ROOT / "examples" / "configs" / "config.json"
OUTPUT_PATH = ROOT / "outputs" / "presentations" / "ttft_assessment_onepager.pptx"


def _add_textbox(slide, left, top, width, height, text, font_size=20, bold=False,
                 color=RGBColor(30, 41, 59), align=PP_ALIGN.LEFT, name=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    if name:
        box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = "Aptos"
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    return box


def _add_panel(slide, left, top, width, height, fill_rgb, line_rgb=None, radius=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius is None else MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = line_rgb or fill_rgb
    shape.line.width = Pt(1)
    return shape


def _add_bullet_list(slide, left, top, width, height, items, font_size=18, color=RGBColor(51, 65, 85)):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        font = p.runs[0].font
        font.name = "Aptos"
        font.size = Pt(font_size)
        font.color.rgb = color
    return box


def main() -> None:
    config = json.loads(CONFIG_PATH.read_text(encoding="utf-8-sig"))
    est_seq8 = build_estimate(config, sequence_length=8)
    est_seq512 = build_estimate(config, sequence_length=512)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(246, 248, 252)

    # Header band
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(15, 23, 42)
    header.line.color.rgb = RGBColor(15, 23, 42)

    _add_textbox(
        slide,
        Inches(0.45),
        Inches(0.18),
        Inches(7.8),
        Inches(0.36),
        "DeepSeek 1.5B FP16 的 TTFT 评估方法",
        font_size=28,
        bold=True,
        color=RGBColor(255, 255, 255),
    )
    _add_textbox(
        slide,
        Inches(0.47),
        Inches(0.53),
        Inches(6.8),
        Inches(0.2),
        "基于主瓶颈 GEMM 的一阶近似，用于当前硬件设计阶段的快速估算",
        font_size=12,
        color=RGBColor(191, 219, 254),
    )

    # Top-left: objective
    _add_panel(slide, Inches(0.35), Inches(1.1), Inches(3.0), Inches(2.2), RGBColor(255, 255, 255), RGBColor(226, 232, 240))
    _add_textbox(slide, Inches(0.55), Inches(1.25), Inches(1.6), Inches(0.3), "1. 评估目标", 20, True)
    _add_bullet_list(
        slide,
        Inches(0.55),
        Inches(1.65),
        Inches(2.5),
        Inches(1.3),
        [
            "目标指标：TTFT（prefill latency）",
            "评估对象：DeepSeek 1.5B FP16",
            "用途：设计阶段快速估算，不替代完整性能模型",
        ],
        font_size=16,
    )

    # Top-middle: config
    _add_panel(slide, Inches(3.55), Inches(1.1), Inches(3.0), Inches(2.2), RGBColor(255, 255, 255), RGBColor(226, 232, 240))
    _add_textbox(slide, Inches(3.75), Inches(1.25), Inches(1.8), Inches(0.3), "2. 基本配置", 20, True)
    _add_bullet_list(
        slide,
        Inches(3.75),
        Inches(1.65),
        Inches(2.55),
        Inches(1.3),
        [
            "hidden = 1536, intermediate = 8960",
            "heads = 12, KV heads = 2, head_dim = 128",
            "28 slices, 4 slices / head, slice freq = 1 GHz",
        ],
        font_size=16,
    )

    # Top-right: execution assumptions
    _add_panel(slide, Inches(6.75), Inches(1.1), Inches(6.2), Inches(2.2), RGBColor(255, 255, 255), RGBColor(226, 232, 240))
    _add_textbox(slide, Inches(6.95), Inches(1.25), Inches(2.2), Inches(0.3), "3. 执行侧近似", 20, True)
    _add_bullet_list(
        slide,
        Inches(6.95),
        Inches(1.65),
        Inches(5.6),
        Inches(1.3),
        [
            "12 个 attention heads 无法填满 7 个 cluster，因此按 14 heads 做 execution padding",
            "执行侧 attention hidden = 14 x 128 = 1792；FFN 仍保持 intermediate = 8960",
            "K/V 需要各 cluster 参与完整计算，但延迟口径按 cluster 并行处理",
        ],
        font_size=16,
    )

    # Bottom-left: method
    _add_panel(slide, Inches(0.35), Inches(3.55), Inches(6.1), Inches(3.15), RGBColor(255, 255, 255), RGBColor(226, 232, 240))
    _add_textbox(slide, Inches(0.55), Inches(3.7), Inches(1.9), Inches(0.3), "4. 估算方法", 20, True)
    _add_bullet_list(
        slide,
        Inches(0.55),
        Inches(4.12),
        Inches(5.6),
        Inches(1.45),
        [
            "主瓶颈按 9 个 GEMM 处理：7 个 ring_gemm + qkt / sv 两个 local_gemm",
            "GEMM latency = 2MNK / (0.92 x 256 ops/cycle)",
            "非 GEMM 算子采用有效带宽近似，作为次要项合并估算",
        ],
        font_size=16,
    )

    metric_panel = _add_panel(slide, Inches(0.62), Inches(5.55), Inches(5.55), Inches(0.82), RGBColor(239, 246, 255), RGBColor(147, 197, 253))
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(5.79),
        Inches(5.0),
        Inches(0.25),
        "TTFT 仅用于设计阶段对比，不替代后续实测或完整 analytical model validation",
        14,
        True,
        RGBColor(30, 64, 175),
    )

    # Bottom-right: results
    _add_panel(slide, Inches(6.75), Inches(3.55), Inches(6.2), Inches(3.15), RGBColor(255, 255, 255), RGBColor(226, 232, 240))
    _add_textbox(slide, Inches(6.95), Inches(3.7), Inches(2.5), Inches(0.3), "5. 当前结果", 20, True)

    result_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.05), Inches(4.2), Inches(5.6), Inches(1.25))
    result_band.fill.solid()
    result_band.fill.fore_color.rgb = RGBColor(15, 23, 42)
    result_band.line.color.rgb = RGBColor(15, 23, 42)
    _add_textbox(
        slide,
        Inches(7.4),
        Inches(4.46),
        Inches(5.0),
        Inches(0.6),
        f"seq = 8  -> TTFT ≈ {est_seq8['total_ms']:.2f} ms\nseq = 512 -> TTFT ≈ {est_seq512['total_ms']:.1f} ms",
        24,
        True,
        RGBColor(255, 255, 255),
    )

    _add_bullet_list(
        slide,
        Inches(6.98),
        Inches(5.75),
        Inches(5.5),
        Inches(0.65),
        [
            "TTFT 会随模型规模、输入长度、精度、并行策略和 runtime 实现显著变化",
            "因此汇报时建议同时注明模型名、precision、input tokens 和近似假设",
        ],
        font_size=15,
        color=RGBColor(71, 85, 105),
    )

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(OUTPUT_PATH)
        print(OUTPUT_PATH)
    except PermissionError:
        fallback = OUTPUT_PATH.with_name("ttft_assessment_onepager_updated.pptx")
        prs.save(fallback)
        print(fallback)


if __name__ == "__main__":
    main()

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from address_remapping.hardware import HardwareSpec
from address_remapping.model_parser import expand_model_spec
from address_remapping.performance import _ring_gemm_execution_geometry, analyze_graph_performance
from address_remapping.roofline import build_roofline_summary
from address_remapping.rmsnorm_bridge import normalize_graph_spec


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BLUE = RGBColor(36, 98, 210)
LIGHT_BLUE = RGBColor(230, 239, 255)
ORANGE = RGBColor(235, 141, 53)
LIGHT_ORANGE = RGBColor(255, 239, 223)
RED = RGBColor(225, 79, 79)
LIGHT_RED = RGBColor(255, 232, 230)
PURPLE = RGBColor(127, 103, 191)
LIGHT_PURPLE = RGBColor(239, 233, 250)
GREEN = RGBColor(55, 153, 103)
LIGHT_GREEN = RGBColor(232, 247, 239)
GRAY = RGBColor(116, 123, 135)
LIGHT_GRAY = RGBColor(244, 246, 249)
DARK = RGBColor(33, 39, 49)
WHITE = RGBColor(255, 255, 255)


def add_textbox(slide, x, y, w, h, text, *, font_size=18, bold=False, color=DARK,
                fill=None, line=None, align=PP_ALIGN.LEFT, font_name="Aptos", margin=0.07):
    shape = slide.shapes.add_textbox(x, y, w, h)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.0)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return shape


def add_box(slide, x, y, w, h, text, *, fill, line, font_size=16, bold=False, color=DARK,
            align=PP_ALIGN.CENTER, shape_type=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.35)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Aptos"
    return shape


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.55), Inches(0.22), Inches(11.5), Inches(0.6), title, font_size=28, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.58), Inches(0.78), Inches(12.0), Inches(0.32), subtitle, font_size=11, color=GRAY)


def add_footer(slide, text):
    add_textbox(slide, Inches(0.55), Inches(7.05), Inches(12.1), Inches(0.2), text, font_size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def add_arrow(slide, x1, y1, x2, y2, *, color=GRAY, width=2.0):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_lane(slide, label, y, *, color, sublabel=None):
    add_textbox(slide, Inches(0.1), y - Inches(0.03), Inches(1.35), Inches(0.45), label, font_size=17, bold=True, color=color)
    if sublabel:
        add_textbox(slide, Inches(0.15), y + Inches(0.26), Inches(1.3), Inches(0.28), sublabel, font_size=10, color=DARK)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.48), y + Inches(0.22), Inches(12.6), y + Inches(0.22))
    line.line.color.rgb = GRAY
    line.line.width = Pt(1.0)
    return line


def add_event_sequence(slide, x_start, y, labels, *, fill, line, w=0.44, gap=0.06, font_size=10):
    x = x_start
    for label in labels:
        add_box(slide, x, y, Inches(w), Inches(0.36), label, fill=fill, line=line, font_size=font_size, bold=False)
        x += Inches(w + gap)


def add_time_axis(slide, x0, y, ticks, *, scale=0.5):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x0, y, x0 + Inches(scale * ticks), y)
    line.line.color.rgb = DARK
    line.line.width = Pt(1.2)
    line.line.end_arrowhead = True
    for i in range(ticks + 1):
        x = x0 + Inches(scale * i)
        tick = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y - Inches(0.07), x, y + Inches(0.07))
        tick.line.color.rgb = DARK
        tick.line.width = Pt(1.0)
        add_textbox(slide, x - Inches(0.08), y - Inches(0.26), Inches(0.16), Inches(0.16), str(i), font_size=8, align=PP_ALIGN.CENTER)


def fmt_list(values):
    out = []
    for value in values:
        fv = float(value)
        out.append(str(int(fv)) if fv.is_integer() else f"{fv:.1f}")
    return "[" + ", ".join(out) + "]"


def fmt_percent(value):
    if value is None:
        return "N/A"
    return f"{float(value) * 100.0:.2f}%"


def fmt_percent_delta(simulated, measured):
    if simulated is None or measured in (None, 0):
        return "N/A"
    return f"{((float(simulated) - float(measured)) / float(measured)) * 100.0:.2f}%"


def build_context():
    graph_path = Path("examples/graphs/ring_gemm/ring_gemm_bias.json")
    graph = json.loads(graph_path.read_text(encoding="utf-8"))
    normalized = normalize_graph_spec(graph, require_base_addrs=True)
    expanded = expand_model_spec(normalized) if "model" in normalized else dict(normalized)
    op_name, op_data = next(iter(dict(expanded["ops"]).items()))
    hw = HardwareSpec()
    geom = _ring_gemm_execution_geometry(op_data, hw)
    report = analyze_graph_performance(graph, hw, include_request_traces=False)
    roofline_summary = build_roofline_summary(report, mode="baseline", graph_name=graph_path.stem)
    roof_op = roofline_summary["operators"][0]
    op_report = report["modes"]["baseline"]["op_breakdown"][0]
    output_port = next(iter(op_data["outputs"]))
    streams = op_report["streams"]
    a_requests = sum(stream["request_count"] for stream in streams if stream["role"] == "A")
    b_requests = sum(stream["request_count"] for stream in streams if stream["role"] == "B")
    write_requests = sum(stream["request_count"] for stream in streams if stream["role"] == "writeback")
    ring_t = op_report["ring_microtile_timeline"]
    analytical = op_report["analytical_model"]
    return {
        "graph": graph,
        "op_name": op_name,
        "op_data": op_data,
        "geometry": geom,
        "op_report": op_report,
        "a_shape": op_data["inputs"]["inA"]["resolved_shape"],
        "b_shape": op_data["inputs"]["inB"]["resolved_shape"],
        "out_shape": op_data["outputs"][output_port]["resolved_shape"],
        "streams": streams,
        "a_requests": a_requests,
        "b_requests": b_requests,
        "write_requests": write_requests,
        "ring_t": ring_t,
        "analytical": analytical,
        "roofline_summary": roofline_summary,
        "roof_op": roof_op,
    }


def slide_summary_table(prs, ctx):
    roof = ctx["roof_op"]
    analytical_compute = fmt_percent(roof.get("analytical_efficiency"))
    analytical_bandwidth = fmt_percent(roof.get("analytical_bandwidth_utilization"))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "ring_gemm 单算子性能表", "指标口径对齐 export_op_performance_summary.py / roofline compact CSV")
    add_box(slide, Inches(0.72), Inches(1.12), Inches(11.9), Inches(0.78),
            "对象：ring_gemm_fp16_fp16_fp16_0  |  mode=baseline  |  running example=ring_gemm_bias.json",
            fill=LIGHT_GREEN, line=GREEN, font_size=16, bold=True)

    rows = [
        ("latency_cycles", f"{int(roof['latency_cycles'])}"),
        ("hardware_measured_cycles", f"{int(roof['hardware_measured_cycles'])}" if roof.get("hardware_measured_cycles") is not None else "N/A"),
        ("simulator_error_vs_measured(%)", fmt_percent_delta(roof.get("latency_cycles"), roof.get("hardware_measured_cycles"))),
        ("roofline_compute_utilization(%)", fmt_percent(roof.get("roofline_compute_utilization"))),
        ("analytical_compute_utilization(%)", fmt_percent(roof.get("analytical_efficiency"))),
        ("measured_compute_utilization(%)", fmt_percent(roof.get("measured_efficiency"))),
        ("roofline_bandwidth_utilization(%)", fmt_percent(roof.get("roofline_bandwidth_utilization"))),
        ("analytical_bandwidth_utilization(%)", fmt_percent(roof.get("analytical_bandwidth_utilization"))),
        ("measured_bandwidth_utilization(%)", fmt_percent(roof.get("measured_bandwidth_utilization"))),
    ]
    x0 = Inches(0.86)
    y0 = Inches(2.08)
    c0 = Inches(6.3)
    c1 = Inches(2.45)
    row_h = Inches(0.42)
    add_box(slide, x0, y0, c0, row_h, "metric", fill=LIGHT_BLUE, line=BLUE, font_size=14, bold=True)
    add_box(slide, x0 + c0, y0, c1, row_h, "value", fill=LIGHT_BLUE, line=BLUE, font_size=14, bold=True)
    for idx, (metric, value) in enumerate(rows, start=1):
        y = y0 + row_h * idx
        fill = WHITE if idx % 2 else LIGHT_GRAY
        add_box(slide, x0, y, c0, row_h, metric, fill=fill, line=GRAY, font_size=13, align=PP_ALIGN.LEFT)
        add_box(slide, x0 + c0, y, c1, row_h, value, fill=fill, line=GRAY, font_size=13)

    add_box(slide, Inches(9.85), Inches(2.08), Inches(2.18), Inches(0.88), f"work_ops\n{int(roof['work_ops'])}", fill=LIGHT_BLUE, line=BLUE, font_size=16, bold=True)
    add_box(slide, Inches(9.85), Inches(3.18), Inches(2.18), Inches(0.88), f"total_bytes\n{int(roof['total_bytes'])}", fill=LIGHT_ORANGE, line=ORANGE, font_size=16, bold=True)
    add_box(slide, Inches(9.85), Inches(4.28), Inches(2.18), Inches(0.88), f"AI (ops/B)\n{roof['arithmetic_intensity_ops_per_byte']:.2f}", fill=LIGHT_PURPLE, line=PURPLE, font_size=16, bold=True)
    add_box(slide, Inches(9.35), Inches(5.55), Inches(2.95), Inches(1.0),
            "当前结论\nsimulator 比 measured 慢约 7.42%\ncompute / bandwidth 利用率都还低于 roofline 上限",
            fill=LIGHT_RED, line=RED, font_size=14, bold=True)
    add_footer(slide, "Slide 1 · 单算子性能表")


def slide_title(prs, ctx):
    op = ctx["op_report"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "ring_gemm 性能分析方法图", "loop nest → template tile → load event → microtile timeline → final bounds")
    add_box(slide, Inches(0.65), Inches(1.2), Inches(3.25), Inches(1.25),
            "运行样例\nring_gemm_bias.json\nGlobal: M=64, N=128, K=16\nLocal: A=64×4, B=16×32, Out=64×32",
            fill=LIGHT_BLUE, line=BLUE, font_size=18, bold=True)
    add_box(slide, Inches(4.2), Inches(1.2), Inches(2.2), Inches(1.25),
            "事件规模\nA events = 4\nB events = 16",
            fill=LIGHT_ORANGE, line=ORANGE, font_size=18, bold=True)
    add_box(slide, Inches(6.7), Inches(1.2), Inches(2.4), Inches(1.25),
            "关键瓶颈\nmemory = 1100\nring = 612",
            fill=LIGHT_RED, line=RED, font_size=18, bold=True)
    add_box(slide, Inches(9.4), Inches(1.2), Inches(2.5), Inches(1.25),
            f"最终延迟\nlatency = {int(op['latency_cycles'])}",
            fill=LIGHT_GREEN, line=GREEN, font_size=18, bold=True)

    add_box(slide, Inches(0.7), Inches(3.0), Inches(11.6), Inches(2.95),
            "", fill=WHITE, line=GRAY, font_size=14)
    add_textbox(slide, Inches(0.95), Inches(3.15), Inches(11.1), Inches(0.3),
                "方法总览：不是直接对 footprint 求 max(compute, memory, AG)，而是先生成事件，再让 A/B/ring/PE/psum/writeback 在时间线上相互约束。", font_size=16, bold=True)
    add_box(slide, Inches(1.0), Inches(4.0), Inches(1.45), Inches(0.8), "loop nest", fill=LIGHT_BLUE, line=BLUE, font_size=17, bold=True)
    add_box(slide, Inches(2.75), Inches(4.0), Inches(1.75), Inches(0.8), "template tiles", fill=LIGHT_BLUE, line=BLUE, font_size=17, bold=True)
    add_box(slide, Inches(4.85), Inches(4.0), Inches(1.75), Inches(0.8), "load events", fill=LIGHT_ORANGE, line=ORANGE, font_size=17, bold=True)
    add_box(slide, Inches(7.0), Inches(4.0), Inches(2.0), Inches(0.8), "microtile timeline", fill=LIGHT_RED, line=RED, font_size=17, bold=True)
    add_box(slide, Inches(9.45), Inches(4.0), Inches(1.85), Inches(0.8), "final bounds", fill=LIGHT_GREEN, line=GREEN, font_size=17, bold=True)
    add_arrow(slide, Inches(2.45), Inches(4.4), Inches(2.75), Inches(4.4), color=GRAY)
    add_arrow(slide, Inches(4.5), Inches(4.4), Inches(4.85), Inches(4.4), color=GRAY)
    add_arrow(slide, Inches(6.6), Inches(4.4), Inches(7.0), Inches(4.4), color=GRAY)
    add_arrow(slide, Inches(9.05), Inches(4.4), Inches(9.45), Inches(4.4), color=GRAY)
    add_textbox(slide, Inches(1.0), Inches(5.15), Inches(10.6), Inches(0.4),
                "核心要点：B 从 64 到 128 requests 的增长来自 loop-scope 重复读取；A/B ready 现在都以 full-group completion 为准。", font_size=18, color=DARK, align=PP_ALIGN.CENTER)
    add_footer(slide, "Slide 2 · 总览")


def slide_loop_scope(prs, ctx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Loop Nest 抽象模型：事件从哪里来")
    add_box(slide, Inches(0.65), Inches(1.2), Inches(4.1), Inches(5.8), "", fill=LIGHT_GRAY, line=GRAY, font_size=14)
    add_textbox(slide, Inches(0.95), Inches(1.45), Inches(3.55), Inches(4.9),
                "for M_0 in range(2)\n"
                "  for N_0 in range(1)\n"
                "    for KM_0 in range(2)\n"
                "      ag0(...)      # A loads\n"
                "    for KN_0 in range(4)\n"
                "      ag1(...)\n"
                "      ag2(...)      # B loads\n"
                "    ag4(...)        # writeback",
                font_size=22, font_name="Consolas", color=DARK)

    add_box(slide, Inches(5.2), Inches(1.45), Inches(2.65), Inches(1.3),
            "A load scope\n(M_0, N_0, KM_0)\n2 × 1 × 2 = 4 events",
            fill=LIGHT_BLUE, line=BLUE, font_size=18, bold=True)
    add_box(slide, Inches(5.2), Inches(3.15), Inches(2.65), Inches(1.3),
            "B load scope\n(M_0, N_0, k_tile)\n2 × 1 × 8 = 16 events",
            fill=LIGHT_ORANGE, line=ORANGE, font_size=18, bold=True)
    add_box(slide, Inches(5.2), Inches(4.85), Inches(2.65), Inches(1.3),
            "writeback scope\noutput tile complete\n2 write groups",
            fill=LIGHT_RED, line=RED, font_size=18, bold=True)

    add_box(slide, Inches(8.25), Inches(1.4), Inches(4.1), Inches(4.95),
            "抽象语义\n\n"
            "• A 的事件重复由 (M_0, N_0, KM_0) 决定\n"
            "• B 的事件重复由 (M_0, N_0, KN_0 / k_tile) 决定\n"
            "• 因为 M_0 前进时会重新进入 B 的 load scope，B footprint=64 req 会扩张成 execution-realistic 128 req\n"
            "• 这个 loop tree 才是 request multiplicity 的权威来源，不是 tensor footprint 本身",
            fill=WHITE, line=GRAY, font_size=18, align=PP_ALIGN.LEFT)
    add_footer(slide, "Slide 3 · loop scope 决定 event multiplicity")


def slide_template_to_events(prs, ctx):
    g = ctx["geometry"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Template Tile → Load Event：几何唯一块如何变成真实事件")

    add_box(slide, Inches(0.55), Inches(1.2), Inches(3.0), Inches(1.1),
            f"footprint requests\nA={ctx['a_requests']}  B=64  W={ctx['write_requests']}",
            fill=LIGHT_BLUE, line=BLUE, font_size=18, bold=True)
    add_box(slide, Inches(3.9), Inches(1.2), Inches(2.4), Inches(1.1),
            "_split_requests_exact()", fill=LIGHT_GRAY, line=GRAY, font_size=18, bold=True)
    add_box(slide, Inches(6.65), Inches(1.2), Inches(2.9), Inches(1.1),
            "template tiles\nA=4, B=8, Out=2", fill=LIGHT_BLUE, line=BLUE, font_size=18, bold=True)
    add_box(slide, Inches(9.95), Inches(1.2), Inches(2.4), Inches(1.1),
            "_expand_ring_gemm_load_events()", fill=LIGHT_GRAY, line=GRAY, font_size=16, bold=True)
    add_box(slide, Inches(9.55), Inches(2.7), Inches(2.8), Inches(1.2),
            f"load events\nA={ctx['op_report']['a_load_event_count']}\nB={ctx['op_report']['b_load_event_count']}",
            fill=LIGHT_ORANGE, line=ORANGE, font_size=18, bold=True)
    add_arrow(slide, Inches(3.55), Inches(1.75), Inches(3.9), Inches(1.75))
    add_arrow(slide, Inches(6.3), Inches(1.75), Inches(6.65), Inches(1.75))
    add_arrow(slide, Inches(9.55), Inches(1.75), Inches(9.95), Inches(1.75))

    add_box(slide, Inches(0.8), Inches(3.0), Inches(3.2), Inches(2.8),
            "A template tiles\n"
            "T0=(M[0:32],K[0:2])\n"
            "T1=(M[0:32],K[2:4])\n"
            "T2=(M[32:64],K[0:2])\n"
            "T3=(M[32:64],K[2:4])",
            fill=WHITE, line=BLUE, font_size=18, align=PP_ALIGN.LEFT)
    add_box(slide, Inches(4.35), Inches(3.0), Inches(3.4), Inches(2.8),
            "B template tiles\n"
            "8 unique K-slices over N[0:32]\n"
            "each tile = 2×32 = 128B = 8 req\n\n"
            "几何上只有 8 个 unique B tiles",
            fill=WHITE, line=ORANGE, font_size=18, align=PP_ALIGN.LEFT)
    add_box(slide, Inches(8.15), Inches(3.0), Inches(4.1), Inches(2.8),
            "event semantics\n"
            "A: 4 template tiles → 4 load events\n"
            "B: 8 template tiles → 16 load events\n\n"
            "因此：\n"
            "b_requests=64 只是 footprint\n"
            f"expanded_b_requests={ctx['b_requests']} 才是执行真实请求流",
            fill=LIGHT_GREEN, line=GREEN, font_size=18, align=PP_ALIGN.LEFT)
    add_footer(slide, "Slide 4 · template tile 与 load event 是两个层次")


def slide_event_ready_bank_timeline(prs, ctx):
    rt = ctx["ring_t"]
    a_tiles = rt["a_buffer_timeline"]["tiles"]
    b_tiles = rt["b_buffer_timeline"]["tiles"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "A/B event ready 如何得到：ring_gemm 中的 bank timeline 建模")

    add_box(slide, Inches(0.7), Inches(1.2), Inches(2.05), Inches(0.9),
            "A/B load event group\n(one event = many requests)", fill=LIGHT_BLUE, line=BLUE, font_size=18, bold=True)
    add_box(slide, Inches(3.0), Inches(1.2), Inches(2.05), Inches(0.9),
            "future_reads /\nfuture_writes", fill=LIGHT_ORANGE, line=ORANGE, font_size=18, bold=True)
    add_box(slide, Inches(5.3), Inches(1.2), Inches(2.05), Inches(0.9),
            "ready queues\nby bank", fill=LIGHT_RED, line=RED, font_size=18, bold=True)
    add_box(slide, Inches(7.6), Inches(1.2), Inches(2.05), Inches(0.9),
            "per-bank service\n(row / phase / drain)", fill=LIGHT_PURPLE, line=PURPLE, font_size=18, bold=True)
    add_box(slide, Inches(9.9), Inches(1.2), Inches(2.05), Inches(0.9),
            "group ready\n(full-group completion)", fill=LIGHT_GREEN, line=GREEN, font_size=18, bold=True)

    add_arrow(slide, Inches(2.75), Inches(1.65), Inches(3.0), Inches(1.65))
    add_arrow(slide, Inches(5.05), Inches(1.65), Inches(5.3), Inches(1.65))
    add_arrow(slide, Inches(7.35), Inches(1.65), Inches(7.6), Inches(1.65))
    add_arrow(slide, Inches(9.65), Inches(1.65), Inches(9.9), Inches(1.65))

    add_box(slide, Inches(0.9), Inches(2.5), Inches(2.3), Inches(1.45),
            "A event 0\n8 requests\nrelease = 0\nready = 28", fill=WHITE, line=BLUE, font_size=17, bold=True)
    add_box(slide, Inches(0.9), Inches(4.25), Inches(2.3), Inches(1.45),
            "B event 0\n8 requests\nrelease = 0\nready = 88", fill=WHITE, line=ORANGE, font_size=17, bold=True)

    add_box(slide, Inches(3.45), Inches(2.4), Inches(3.8), Inches(3.55),
            "bank timeline 抽象过程\n\n"
            "1. event 的 requests 先挂到 future_reads / future_writes\n"
            "2. release_cycle 到了之后，按 bank_id 进入 ready queues\n"
            "3. per-bank timeline 决定每条 request 何时完成\n"
            "4. 只有当同一 group 的全部 requests 都完成时，\n"
            "   才写 completion_by_group[group_key]\n"
            "5. 这个完成时间就是 event ready time",
            fill=LIGHT_GRAY, line=GRAY, font_size=17, align=PP_ALIGN.LEFT)

    add_box(slide, Inches(7.7), Inches(2.55), Inches(4.2), Inches(1.5),
            "修正后的关键语义\n\n"
            "event ready ≠ first request done\n"
            "event ready = all requests in the group done",
            fill=LIGHT_RED, line=RED, font_size=18, bold=True, align=PP_ALIGN.LEFT)
    add_box(slide, Inches(7.7), Inches(4.4), Inches(4.2), Inches(1.5),
            f"当前例子\n\nA0 ready = {int(a_tiles[0]['load_ready_cycle'])}, A1 ready = {int(a_tiles[1]['load_ready_cycle'])}\n"
            f"B0 ready = {int(b_tiles[0]['load_ready_cycle'])}, B1 ready = {int(b_tiles[1]['load_ready_cycle'])}",
            fill=WHITE, line=GREEN, font_size=18, bold=True, align=PP_ALIGN.LEFT)

    add_box(slide, Inches(0.8), Inches(6.15), Inches(11.45), Inches(0.65),
            "这一页的核心：A/B event ready 不是公式直接给定，而是 bank timeline 对整组 requests 做 full-group completion 后反馈给上层 microtile timeline 的时间。",
            fill=LIGHT_GREEN, line=GREEN, font_size=16, bold=True, align=PP_ALIGN.LEFT)
    add_footer(slide, "Slide 5 · event ready 来自 bank timeline，而不是单条 request")


def slide_output_tile_event_graph(prs, ctx):
    rt = ctx["ring_t"]
    a_tiles = rt["a_buffer_timeline"]["tiles"]
    b_ready = [int(v) for v in rt["b_ready_cycles"][:8]]
    ring_ready = [int(v) for v in rt["ring_ready_cycles"][:8]]
    compute_start = [int(v) for v in rt["compute_start_cycles"][:8]]
    compute_end = [int(v) for v in rt["compute_end_cycles"][:8]]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "一个 output tile 内部：A/B event ready 与 coarse compute 时间")

    add_box(slide, Inches(0.72), Inches(1.08), Inches(2.35), Inches(0.82),
            "output tile 0\nM[0:32], N[0:32]", fill=LIGHT_PURPLE, line=PURPLE, font_size=18, bold=True)
    add_textbox(slide, Inches(3.15), Inches(1.15), Inches(8.5), Inches(0.26),
                "注意：A/B event 这里只标 ready；只有 coarse compute 才标 start/end。", font_size=12, color=GRAY)

    x0 = Inches(3.1)
    lane_w = Inches(8.7)
    add_lane(slide, "1) A event", Inches(2.0), color=BLUE, sublabel="local-ready + ring-ready")
    add_lane(slide, "2) B event", Inches(3.35), color=ORANGE, sublabel="B ready")
    add_lane(slide, "3) coarse compute", Inches(4.95), color=RED, sublabel="start / end")

    a0_x = x0
    a1_x = x0 + Inches(4.35)
    add_box(slide, a0_x, Inches(1.82), Inches(1.85), Inches(0.4), "A event 0", fill=LIGHT_BLUE, line=BLUE, font_size=14, bold=True)
    add_box(slide, a0_x + Inches(2.0), Inches(1.82), Inches(2.05), Inches(0.4), "ring-ready → coarse k=0..3", fill=WHITE, line=BLUE, font_size=13, bold=True)
    add_textbox(slide, a0_x, Inches(2.25), Inches(4.15), Inches(0.32),
                f"local ready {int(a_tiles[0]['load_ready_cycle'])} | ring-ready {fmt_list(ring_ready[:4])}",
                font_size=11, color=BLUE, align=PP_ALIGN.CENTER)

    add_box(slide, a1_x, Inches(1.82), Inches(1.85), Inches(0.4), "A event 1", fill=LIGHT_BLUE, line=BLUE, font_size=14, bold=True)
    add_box(slide, a1_x + Inches(2.0), Inches(1.82), Inches(2.05), Inches(0.4), "ring-ready → coarse k=4..7", fill=WHITE, line=BLUE, font_size=13, bold=True)
    add_textbox(slide, a1_x, Inches(2.25), Inches(4.15), Inches(0.32),
                f"local ready {int(a_tiles[1]['load_ready_cycle'])} | ring-ready {fmt_list(ring_ready[4:8])}",
                font_size=11, color=BLUE, align=PP_ALIGN.CENTER)

    for idx in range(8):
        x = x0 + Inches(idx * 1.05)
        add_box(slide, x, Inches(3.18), Inches(0.82), Inches(0.42), f"B{idx}", fill=LIGHT_ORANGE, line=ORANGE, font_size=13, bold=True)
        add_textbox(slide, x - Inches(0.02), Inches(3.62), Inches(0.86), Inches(0.24), f"ready {b_ready[idx]}", font_size=10, color=ORANGE, align=PP_ALIGN.CENTER)

        add_box(slide, x - Inches(0.02), Inches(4.72), Inches(0.86), Inches(0.56), f"coarse\nk={idx}", fill=LIGHT_RED, line=RED, font_size=12, bold=True)
        add_textbox(slide, x - Inches(0.05), Inches(5.33), Inches(0.92), Inches(0.22), f"start {compute_start[idx]}", font_size=10, color=RED, align=PP_ALIGN.CENTER)
        add_textbox(slide, x - Inches(0.05), Inches(5.56), Inches(0.92), Inches(0.22), f"end {compute_end[idx]}", font_size=10, color=RED, align=PP_ALIGN.CENTER)
        add_arrow(slide, x + Inches(0.4), Inches(3.6), x + Inches(0.4), Inches(4.72), color=ORANGE, width=1.3)

    add_arrow(slide, a0_x + Inches(4.0), Inches(2.02), x0 + Inches(0.4), Inches(4.95), color=BLUE, width=1.2)
    add_arrow(slide, a0_x + Inches(4.0), Inches(2.02), x0 + Inches(3.55), Inches(4.95), color=BLUE, width=1.2)
    add_arrow(slide, a1_x + Inches(4.0), Inches(2.02), x0 + Inches(4.6), Inches(4.95), color=BLUE, width=1.2)
    add_arrow(slide, a1_x + Inches(4.0), Inches(2.02), x0 + Inches(7.75), Inches(4.95), color=BLUE, width=1.2)

    add_box(slide, Inches(0.78), Inches(6.32), Inches(11.45), Inches(0.58),
            "coarse 粒度：A(32×2) × B(2×32) 对一个 32×32 output tile 的单个 K-slice 做累加。A/B event 只负责给出 ready；coarse_start 才是 compute 真正开始的时间。",
            fill=LIGHT_GREEN, line=GREEN, font_size=15, bold=True, align=PP_ALIGN.LEFT)
    add_footer(slide, "Slide 6 · A/B event ready 与 coarse compute 时间要分开看")


def slide_resource_dependency(prs, ctx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "资源约束图：ping/pong、ring-ready、PE、psum 如何共同决定 coarse_start")
    add_box(slide, Inches(0.7), Inches(1.35), Inches(2.1), Inches(0.95), "A ping / pong", fill=LIGHT_BLUE, line=BLUE, font_size=18, bold=True)
    add_box(slide, Inches(0.7), Inches(2.65), Inches(2.1), Inches(0.95), "B ping / pong", fill=LIGHT_ORANGE, line=ORANGE, font_size=18, bold=True)
    add_box(slide, Inches(0.7), Inches(3.95), Inches(2.1), Inches(0.95), "psum ping / pong", fill=LIGHT_PURPLE, line=PURPLE, font_size=18, bold=True)
    add_box(slide, Inches(3.45), Inches(1.35), Inches(2.2), Inches(0.95), "A local ready", fill=WHITE, line=BLUE, font_size=18, bold=True)
    add_box(slide, Inches(3.45), Inches(2.65), Inches(2.2), Inches(0.95), "B ready", fill=WHITE, line=ORANGE, font_size=18, bold=True)
    add_box(slide, Inches(3.45), Inches(3.95), Inches(2.2), Inches(0.95), "PE available", fill=WHITE, line=RED, font_size=18, bold=True)
    add_box(slide, Inches(6.15), Inches(1.35), Inches(2.3), Inches(0.95), "ring_ready\n= A ready + hop", fill=LIGHT_BLUE, line=BLUE, font_size=18, bold=True)
    add_box(slide, Inches(6.15), Inches(3.95), Inches(2.3), Inches(0.95), "psum slot ready", fill=LIGHT_PURPLE, line=PURPLE, font_size=18, bold=True)
    add_box(slide, Inches(8.95), Inches(2.55), Inches(2.35), Inches(1.15),
            "coarse_start\n= max(PE, ring, B, psum)", fill=LIGHT_RED, line=RED, font_size=20, bold=True)
    add_box(slide, Inches(11.55), Inches(2.55), Inches(1.0), Inches(1.15),
            "PE\ncompute", fill=LIGHT_ORANGE, line=ORANGE, font_size=17, bold=True)
    for y1, y2 in [(1.82, 1.82), (3.12, 3.12), (4.42, 4.42)]:
        add_arrow(slide, Inches(2.8), Inches(y1), Inches(3.45), Inches(y2), color=GRAY)
    add_arrow(slide, Inches(5.65), Inches(1.82), Inches(6.15), Inches(1.82), color=GRAY)
    add_arrow(slide, Inches(5.65), Inches(3.12), Inches(8.95), Inches(3.12), color=GRAY)
    add_arrow(slide, Inches(5.65), Inches(4.42), Inches(8.95), Inches(3.12), color=GRAY)
    add_arrow(slide, Inches(8.45), Inches(1.82), Inches(8.95), Inches(3.12), color=GRAY)
    add_arrow(slide, Inches(8.45), Inches(4.42), Inches(8.95), Inches(3.12), color=GRAY)
    add_arrow(slide, Inches(11.3), Inches(3.12), Inches(11.55), Inches(3.12), color=GRAY)
    add_textbox(slide, Inches(0.9), Inches(5.45), Inches(11.3), Inches(0.65),
                "当前实现里：A slot 要等最后一个 ring-hop 消费完才释放；B slot 在当前 coarse compute 后即可释放。", font_size=17, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Slide 7 · coarse_start 由多资源依赖共同决定")


def slide_timeline_figure(prs, ctx):
    rt = ctx["ring_t"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "ring_gemm 抽象执行时间线示意图")
    add_textbox(slide, Inches(0.75), Inches(1.06), Inches(4.0), Inches(0.24),
                "示意风格图：强调事件顺序与依赖，不按真实 cycle 精确缩放。", font_size=11, color=GRAY)
    x0 = Inches(1.7)
    axis_y = Inches(1.55)
    add_time_axis(slide, x0, axis_y, 16, scale=0.62)
    add_textbox(slide, Inches(0.2), Inches(1.37), Inches(1.35), Inches(0.22), "event steps", font_size=11, color=GRAY, align=PP_ALIGN.RIGHT)
    add_lane(slide, "1) A local load", Inches(2.0), color=BLUE, sublabel="A event")
    add_lane(slide, "2) A ring-ready", Inches(2.95), color=BLUE, sublabel="ring hop")
    add_lane(slide, "3) B load", Inches(3.9), color=ORANGE, sublabel="B event")
    add_lane(slide, "4) coarse compute", Inches(4.85), color=RED, sublabel="K-slice accumulate")
    add_lane(slide, "5) psum / writeback", Inches(5.8), color=PURPLE, sublabel="tile complete")

    add_event_sequence(slide, x0, Inches(1.82), ["A0", "A1", "A2", "A3"], fill=LIGHT_BLUE, line=BLUE, w=0.58, gap=1.35, font_size=12)
    add_event_sequence(slide, x0 + Inches(0.25), Inches(2.77), ["hop0", "hop1", "hop2", "hop3", "hop0", "hop1", "hop2", "hop3"], fill=WHITE, line=BLUE, w=0.6, gap=0.4, font_size=11)
    add_event_sequence(slide, x0, Inches(3.72), [f"B{i}" for i in range(8)], fill=LIGHT_ORANGE, line=ORANGE, w=0.52, gap=0.34, font_size=11)
    add_event_sequence(slide, x0 + Inches(0.25), Inches(4.67), [f"C{i}" for i in range(8)], fill=LIGHT_RED, line=RED, w=0.68, gap=0.35, font_size=11)
    add_event_sequence(slide, x0 + Inches(0.9), Inches(5.62), ["psum0", "wb0", "psum1", "wb1"], fill=LIGHT_PURPLE, line=PURPLE, w=0.82, gap=1.1, font_size=11)

    add_box(slide, Inches(9.35), Inches(1.85), Inches(2.6), Inches(1.25),
            f"first values\nring_ready = {fmt_list(rt['ring_ready_cycles'][:4])}\nB ready = {fmt_list(rt['b_ready_cycles'][:4])}",
            fill=WHITE, line=GRAY, font_size=14, align=PP_ALIGN.LEFT)
    add_box(slide, Inches(9.35), Inches(3.45), Inches(2.6), Inches(1.25),
            f"compute start = {fmt_list(rt['compute_start_cycles'][:4])}\ncompute end = {fmt_list(rt['compute_end_cycles'][:4])}",
            fill=WHITE, line=GRAY, font_size=14, align=PP_ALIGN.LEFT)
    add_box(slide, Inches(8.9), Inches(5.95), Inches(3.35), Inches(0.6),
            "final latency = max(compute_pipeline, ring, memory, AG)", fill=LIGHT_GREEN, line=GREEN, font_size=15, bold=True)
    add_footer(slide, "Slide 8 · 总时间线图：A/B/ring/compute/psum/writeback")


def slide_bank_model(prs, ctx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "bank-level 简化模型：事件如何进入 memory timeline")
    add_box(slide, Inches(0.75), Inches(1.25), Inches(2.1), Inches(0.95), "future_reads / future_writes", fill=LIGHT_BLUE, line=BLUE, font_size=17, bold=True)
    add_box(slide, Inches(3.1), Inches(1.25), Inches(2.0), Inches(0.95), "ready queues by bank", fill=LIGHT_ORANGE, line=ORANGE, font_size=17, bold=True)
    add_box(slide, Inches(5.35), Inches(1.25), Inches(2.0), Inches(0.95), "per-bank state", fill=LIGHT_RED, line=RED, font_size=17, bold=True)
    add_box(slide, Inches(7.6), Inches(1.25), Inches(2.1), Inches(0.95), "completion_by_group", fill=LIGHT_PURPLE, line=PURPLE, font_size=17, bold=True)
    add_box(slide, Inches(9.95), Inches(1.25), Inches(2.1), Inches(0.95), "memory_timeline_cycles", fill=LIGHT_GREEN, line=GREEN, font_size=17, bold=True)
    add_arrow(slide, Inches(2.87), Inches(1.73), Inches(3.1), Inches(1.73))
    add_arrow(slide, Inches(5.12), Inches(1.73), Inches(5.35), Inches(1.73))
    add_arrow(slide, Inches(7.37), Inches(1.73), Inches(7.6), Inches(1.73))
    add_arrow(slide, Inches(9.72), Inches(1.73), Inches(9.95), Inches(1.73))

    add_box(slide, Inches(0.85), Inches(3.0), Inches(5.0), Inches(2.55),
            "当前已准确表达\n\n"
            "• request 按 release_cycle 注入\n"
            "• bank 内 row / phase 状态\n"
            "• write buffer occupancy + forced drain\n"
            "• group ready = full-group completion\n"
            "• A/B/writeback event completion 反馈到上层 timeline",
            fill=WHITE, line=GRAY, font_size=18, align=PP_ALIGN.LEFT)
    add_box(slide, Inches(6.2), Inches(3.0), Inches(5.75), Inches(2.55),
            "当前仍然简化\n\n"
            "• ring_gemm 没完全复用 ordinary operator 的 controller / arbiter 语义\n"
            "• bank 侧是 lighter event loop，不是显式 arbiter1 / arbiter2 竞赛\n"
            "• 但它已经足够提供 A/B ready、writeback completion 和 memory_timeline 的关键路径",
            fill=LIGHT_RED, line=RED, font_size=18, align=PP_ALIGN.LEFT)
    add_footer(slide, "Slide 9 · 当前 memory model 的准确性与边界")


def slide_bounds(prs, ctx):
    a = ctx["analytical"]
    op = ctx["op_report"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "最终瓶颈选择：谁最后结束，谁决定 latency")
    add_box(slide, Inches(0.85), Inches(1.55), Inches(2.0), Inches(1.15), "compute\n256", fill=LIGHT_BLUE, line=BLUE, font_size=24, bold=True)
    add_box(slide, Inches(3.15), Inches(1.55), Inches(2.0), Inches(1.15), "AG issue\n256", fill=LIGHT_ORANGE, line=ORANGE, font_size=24, bold=True)
    add_box(slide, Inches(5.45), Inches(1.55), Inches(2.15), Inches(1.15), f"ring\n{int(a['ring_transfer_bound_cycles'])}", fill=LIGHT_RED, line=RED, font_size=24, bold=True)
    add_box(slide, Inches(7.95), Inches(1.55), Inches(2.45), Inches(1.15), f"memory\n{int(a['memory_access_bound_cycles'])}", fill=LIGHT_PURPLE, line=PURPLE, font_size=24, bold=True)
    add_box(slide, Inches(10.75), Inches(1.55), Inches(1.8), Inches(1.15), f"latency\n{int(op['latency_cycles'])}", fill=LIGHT_GREEN, line=GREEN, font_size=24, bold=True)
    add_arrow(slide, Inches(2.9), Inches(2.1), Inches(10.75), Inches(2.1), color=GRAY)
    add_box(slide, Inches(0.95), Inches(3.35), Inches(11.3), Inches(1.2),
            "当前样例的结论：在修正 group ready = full-group completion 之后，memory timeline = 1100 仍然最大，"
            "因此最终瓶颈仍然是 memory，而不是 ideal compute(256) 或 AG(256)。",
            fill=WHITE, line=GRAY, font_size=19, bold=True, align=PP_ALIGN.LEFT)
    add_box(slide, Inches(1.6), Inches(5.1), Inches(9.9), Inches(0.8),
            "max(compute_pipeline, ring_timeline, memory_timeline, ag_issue)  →  最晚结束者决定 latency",
            fill=LIGHT_GREEN, line=GREEN, font_size=20, bold=True)
    add_footer(slide, "Slide 10 · final bound selection")


def slide_takeaways(prs, ctx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "总结：当前 ring_gemm 性能分析图应该传达什么")
    add_box(slide, Inches(0.8), Inches(1.35), Inches(3.8), Inches(4.9),
            "1) 事件不是附属概念，而是分析主体\n\n"
            "template tile 只描述几何唯一块；真正驱动 timeline 的是 A/B load event、coarse compute event、writeback event。",
            fill=LIGHT_BLUE, line=BLUE, font_size=19, align=PP_ALIGN.LEFT)
    add_box(slide, Inches(4.8), Inches(1.35), Inches(3.8), Inches(4.9),
            "2) loop nest 决定真实请求流\n\n"
            "B 从 64 到 128 requests 不是 tensor 变大，而是 M_0 推进时重新进入 B 的 load scope。",
            fill=LIGHT_ORANGE, line=ORANGE, font_size=19, align=PP_ALIGN.LEFT)
    add_box(slide, Inches(8.8), Inches(1.35), Inches(3.8), Inches(4.9),
            "3) latency 来自多资源约束\n\n"
            "A local-ready、A ring-ready、B ready、PE available、psum slot、memory timeline 会共同决定最终性能。",
            fill=LIGHT_RED, line=RED, font_size=19, align=PP_ALIGN.LEFT)
    add_box(slide, Inches(1.0), Inches(6.05), Inches(11.1), Inches(0.7),
            "下一步最值得做的 refinement：把 ring_gemm 的 bank event loop 进一步与 ordinary operator 的 controller/arbiter 语义统一。",
            fill=LIGHT_GREEN, line=GREEN, font_size=17, bold=True)
    add_footer(slide, "Slide 11 · takeaway")


def slide_summary_table(prs, ctx):
    roof = ctx["roof_op"]
    analytical_compute = fmt_percent(roof.get("analytical_efficiency"))
    analytical_bandwidth = fmt_percent(roof.get("analytical_bandwidth_utilization"))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "ring_gemm 单算子性能表", "指标口径对齐 export_op_performance_summary.py / roofline compact CSV")
    add_box(slide, Inches(0.72), Inches(1.12), Inches(11.9), Inches(0.78),
            "对象：ring_gemm_fp16_fp16_fp16_0  |  mode=baseline  |  running example=ring_gemm/ring_gemm_bias.json",
            fill=LIGHT_GREEN, line=GREEN, font_size=16, bold=True)
    rows = [
        ("latency_cycles", f"{int(roof['latency_cycles'])}"),
        ("hardware_measured_cycles", f"{int(roof['hardware_measured_cycles'])}" if roof.get("hardware_measured_cycles") is not None else "N/A"),
        ("simulator_error_vs_measured(%)", fmt_percent_delta(roof.get("latency_cycles"), roof.get("hardware_measured_cycles"))),
        ("roofline_compute_utilization(%)", fmt_percent(roof.get("roofline_compute_utilization"))),
        ("analytical_compute_utilization(%)", fmt_percent(roof.get("analytical_efficiency"))),
        ("measured_compute_utilization(%)", fmt_percent(roof.get("measured_efficiency"))),
        ("roofline_bandwidth_utilization(%)", fmt_percent(roof.get("roofline_bandwidth_utilization"))),
        ("analytical_bandwidth_utilization(%)", fmt_percent(roof.get("analytical_bandwidth_utilization"))),
        ("measured_bandwidth_utilization(%)", fmt_percent(roof.get("measured_bandwidth_utilization"))),
    ]
    x0 = Inches(0.86)
    y0 = Inches(2.08)
    c0 = Inches(6.3)
    c1 = Inches(2.45)
    row_h = Inches(0.42)
    add_box(slide, x0, y0, c0, row_h, "metric", fill=LIGHT_BLUE, line=BLUE, font_size=14, bold=True)
    add_box(slide, x0 + c0, y0, c1, row_h, "value", fill=LIGHT_BLUE, line=BLUE, font_size=14, bold=True)
    for idx, (metric, value) in enumerate(rows, start=1):
        y = y0 + row_h * idx
        fill = WHITE if idx % 2 else LIGHT_GRAY
        add_box(slide, x0, y, c0, row_h, metric, fill=fill, line=GRAY, font_size=13, align=PP_ALIGN.LEFT)
        add_box(slide, x0 + c0, y, c1, row_h, value, fill=fill, line=GRAY, font_size=13)
    add_box(slide, Inches(9.85), Inches(2.08), Inches(2.18), Inches(0.88), f"work_ops\n{int(roof['work_ops'])}", fill=LIGHT_BLUE, line=BLUE, font_size=16, bold=True)
    add_box(slide, Inches(9.85), Inches(3.18), Inches(2.18), Inches(0.88), f"total_bytes\n{int(roof['total_bytes'])}", fill=LIGHT_ORANGE, line=ORANGE, font_size=16, bold=True)
    add_box(slide, Inches(9.85), Inches(4.28), Inches(2.18), Inches(0.88), f"AI (ops/B)\n{roof['arithmetic_intensity_ops_per_byte']:.2f}", fill=LIGHT_PURPLE, line=PURPLE, font_size=16, bold=True)
    add_box(slide, Inches(9.35), Inches(5.55), Inches(2.95), Inches(1.0),
            f"当前结论\nsimulator 相对 measured 误差 {fmt_percent_delta(roof.get('latency_cycles'), roof.get('hardware_measured_cycles'))}\n"
            f"analytical compute={analytical_compute}\nanalytical bandwidth={analytical_bandwidth}",
            fill=LIGHT_RED, line=RED, font_size=14, bold=True)
    add_footer(slide, "Slide 1 · 单算子性能表")


def slide_title(prs, ctx):
    op = ctx["op_report"]
    analytical = ctx["analytical"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "ring_gemm 性能分析方法图", "loop nest → template tile → load event → microtile timeline → final bounds")
    add_box(slide, Inches(0.65), Inches(1.2), Inches(3.25), Inches(1.25),
            "运行样例\nring_gemm/ring_gemm_bias.json\nGlobal: M=64, N=128, K=16\nLocal: A=64×4, B=16×32, Out=64×32",
            fill=LIGHT_BLUE, line=BLUE, font_size=18, bold=True)
    add_box(slide, Inches(4.2), Inches(1.2), Inches(2.2), Inches(1.25),
            "事件规模\nA events = 4\nB events = 16",
            fill=LIGHT_ORANGE, line=ORANGE, font_size=18, bold=True)
    add_box(slide, Inches(6.7), Inches(1.2), Inches(2.4), Inches(1.25),
            f"关键瓶颈\nmemory = {int(analytical['memory_access_bound_cycles'])}\nring = {int(analytical['ring_ready_bound_cycles'])}",
            fill=LIGHT_RED, line=RED, font_size=18, bold=True)
    add_box(slide, Inches(9.4), Inches(1.2), Inches(2.5), Inches(1.25),
            f"最终延迟\nlatency = {int(op['latency_cycles'])}",
            fill=LIGHT_GREEN, line=GREEN, font_size=18, bold=True)
    add_box(slide, Inches(0.7), Inches(3.0), Inches(11.6), Inches(2.95), "", fill=WHITE, line=GRAY, font_size=14)
    add_textbox(slide, Inches(0.95), Inches(3.15), Inches(11.1), Inches(0.3),
                "方法总览：不是直接对 footprint 求 max(compute, memory, AG)，而是先生成事件，再让 A/B/ring/PE/psum/writeback 在时间线上相互约束。", font_size=16, bold=True)
    add_box(slide, Inches(1.0), Inches(4.0), Inches(1.45), Inches(0.8), "loop nest", fill=LIGHT_BLUE, line=BLUE, font_size=17, bold=True)
    add_box(slide, Inches(2.75), Inches(4.0), Inches(1.75), Inches(0.8), "template tiles", fill=LIGHT_BLUE, line=BLUE, font_size=17, bold=True)
    add_box(slide, Inches(4.85), Inches(4.0), Inches(1.75), Inches(0.8), "load events", fill=LIGHT_ORANGE, line=ORANGE, font_size=17, bold=True)
    add_box(slide, Inches(7.0), Inches(4.0), Inches(2.0), Inches(0.8), "microtile timeline", fill=LIGHT_RED, line=RED, font_size=17, bold=True)
    add_box(slide, Inches(9.45), Inches(4.0), Inches(1.85), Inches(0.8), "final bounds", fill=LIGHT_GREEN, line=GREEN, font_size=17, bold=True)
    add_arrow(slide, Inches(2.45), Inches(4.4), Inches(2.75), Inches(4.4), color=GRAY)
    add_arrow(slide, Inches(4.5), Inches(4.4), Inches(4.85), Inches(4.4), color=GRAY)
    add_arrow(slide, Inches(6.6), Inches(4.4), Inches(7.0), Inches(4.4), color=GRAY)
    add_arrow(slide, Inches(9.05), Inches(4.4), Inches(9.45), Inches(4.4), color=GRAY)
    add_textbox(slide, Inches(1.0), Inches(5.15), Inches(10.6), Inches(0.4),
                "核心要点：B 从 64 到 128 requests 的增长来自 loop-scope 重复读取；A/B ready 现在都以 full-group completion 为准。", font_size=18, color=DARK, align=PP_ALIGN.CENTER)
    add_footer(slide, "Slide 2 · 总览")


def slide_bounds(prs, ctx):
    a = ctx["analytical"]
    op = ctx["op_report"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "最终瓶颈选择：谁最后结束，谁决定 latency")
    add_box(slide, Inches(0.85), Inches(1.55), Inches(2.0), Inches(1.15), "compute\n256", fill=LIGHT_BLUE, line=BLUE, font_size=24, bold=True)
    add_box(slide, Inches(3.15), Inches(1.55), Inches(2.0), Inches(1.15), "AG issue\n256", fill=LIGHT_ORANGE, line=ORANGE, font_size=24, bold=True)
    add_box(slide, Inches(5.45), Inches(1.55), Inches(2.15), Inches(1.15), f"ring\n{int(a['ring_transfer_bound_cycles'])}", fill=LIGHT_RED, line=RED, font_size=24, bold=True)
    add_box(slide, Inches(7.95), Inches(1.55), Inches(2.45), Inches(1.15), f"memory\n{int(a['memory_access_bound_cycles'])}", fill=LIGHT_PURPLE, line=PURPLE, font_size=24, bold=True)
    add_box(slide, Inches(10.75), Inches(1.55), Inches(1.8), Inches(1.15), f"latency\n{int(op['latency_cycles'])}", fill=LIGHT_GREEN, line=GREEN, font_size=24, bold=True)
    add_arrow(slide, Inches(2.9), Inches(2.1), Inches(10.75), Inches(2.1), color=GRAY)
    add_box(slide, Inches(0.95), Inches(3.35), Inches(11.3), Inches(1.2),
            f"当前样例的结论：memory timeline = {int(a['memory_access_bound_cycles'])} 仍然最大，因此最终瓶颈仍然是 memory，而不是 ideal compute({int(a['compute_bound_cycles'])}) 或 AG({int(a['ag_issue_bound_cycles'])})。",
            fill=WHITE, line=GRAY, font_size=19, bold=True, align=PP_ALIGN.LEFT)
    add_box(slide, Inches(1.6), Inches(5.1), Inches(9.9), Inches(0.8),
            "max(compute_pipeline, ring_timeline, memory_timeline, ag_issue)  →  最晚结束者决定 latency",
            fill=LIGHT_GREEN, line=GREEN, font_size=20, bold=True)
    add_footer(slide, "Slide 10 · final bound selection")


def build_deck(output_path: Path, notes_path: Path):
    ctx = build_context()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_summary_table(prs, ctx)
    slide_title(prs, ctx)
    slide_loop_scope(prs, ctx)
    slide_template_to_events(prs, ctx)
    slide_event_ready_bank_timeline(prs, ctx)
    slide_output_tile_event_graph(prs, ctx)
    slide_resource_dependency(prs, ctx)
    slide_timeline_figure(prs, ctx)
    slide_bank_model(prs, ctx)
    slide_bounds(prs, ctx)
    slide_takeaways(prs, ctx)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    notes_path.write_text(
        "# ring_gemm 性能分析方法图 Deck\n\n"
        "- 运行样例：`examples/graphs/ring_gemm_bias.json`\n"
        "- 语言风格：中文叙事 + 英文变量/事件名\n"
        "- 图形风格：抽象示意为主，不按真实 cycle 等比例绘制\n"
        f"- 首页性能表口径：latency={int(ctx['roof_op']['latency_cycles'])}, hardware_measured={int(ctx['roof_op']['hardware_measured_cycles'])}, simulator_error={fmt_percent_delta(ctx['roof_op']['latency_cycles'], ctx['roof_op']['hardware_measured_cycles'])}\n"
        f"- 关键数值：A events={ctx['op_report']['a_load_event_count']}，B events={ctx['op_report']['b_load_event_count']}，B requests={ctx['b_requests']}，latency={int(ctx['op_report']['latency_cycles'])}\n"
        "- 关键语义：A/B load-event ready 采用 full-group completion。\n",
        encoding="utf-8",
    )


def build_deck(output_path: Path, notes_path: Path):
    ctx = build_context()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_summary_table(prs, ctx)
    slide_title(prs, ctx)
    slide_loop_scope(prs, ctx)
    slide_template_to_events(prs, ctx)
    slide_event_ready_bank_timeline(prs, ctx)
    slide_output_tile_event_graph(prs, ctx)
    slide_resource_dependency(prs, ctx)
    slide_timeline_figure(prs, ctx)
    slide_bank_model(prs, ctx)
    slide_bounds(prs, ctx)
    slide_takeaways(prs, ctx)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    notes_path.write_text(
        "# ring_gemm 性能分析方法图 Deck\n\n"
        "- 运行样例：`examples/graphs/ring_gemm/ring_gemm_bias.json`\n"
        "- 语言风格：中文叙事 + 英文变量/事件名\n"
        "- 图形风格：抽象示意为主，不按真实 cycle 等比例绘制\n"
        f"- 首页性能表口径：latency={int(ctx['roof_op']['latency_cycles'])}, hardware_measured={int(ctx['roof_op']['hardware_measured_cycles'])}, simulator_error={fmt_percent_delta(ctx['roof_op']['latency_cycles'], ctx['roof_op']['hardware_measured_cycles'])}\n"
        f"- 关键数值：A events={ctx['op_report']['a_load_event_count']}，B events={ctx['op_report']['b_load_event_count']}，B requests={ctx['b_requests']}，latency={int(ctx['op_report']['latency_cycles'])}\n"
        "- 关键语义：A/B load-event ready 采用 full-group completion。\n",
        encoding="utf-8",
    )


REFERENCE_CASE = {
    "label": "M=64,N=128,K=16",
    "latency_cycles": 780.0,
    "hardware_measured_cycles": 819.0,
    "roofline_compute_utilization": 1.0,
    "analytical_efficiency": 0.3282051282051282,
    "measured_efficiency": 0.3125763125763126,
    "roofline_bandwidth_utilization": 0.8125,
    "analytical_bandwidth_utilization": 0.26666666666666666,
    "measured_bandwidth_utilization": 0.25396825396825395,
}


def build_case_summary(graph_relpath: str, *, label_suffix: str = ""):
    graph_path = Path(graph_relpath)
    graph = json.loads(graph_path.read_text(encoding="utf-8"))
    hw = HardwareSpec()
    report = analyze_graph_performance(graph, hw, include_request_traces=False)
    roofline_summary = build_roofline_summary(report, mode="baseline", graph_name=graph_path.stem)
    roof_op = dict(roofline_summary["operators"][0])
    base_label = f"M={graph['shape_bindings']['M']},N={graph['shape_bindings']['N']},K={graph['shape_bindings']['K']}"
    roof_op["label"] = f"{base_label}{label_suffix}"
    return roof_op


def slide_summary_table(prs, ctx):
    current = dict(ctx["roof_op"])
    current["label"] = f"M={ctx['graph']['shape_bindings']['M']},N={ctx['graph']['shape_bindings']['N']},K={ctx['graph']['shape_bindings']['K']}"
    case_28 = build_case_summary(
        "examples/graphs/ring_gemm/ring_gemm_bias_28slices.json",
        label_suffix=" (28 slices)",
    )
    cases = [REFERENCE_CASE, current, case_28]
    metrics = [
        ("latency_cycles", lambda c: f"{int(c['latency_cycles'])}"),
        ("hardware_measured_cycles", lambda c: f"{int(c['hardware_measured_cycles'])}" if c.get("hardware_measured_cycles") is not None else "N/A"),
        ("simulator_error_vs_measured(%)", lambda c: fmt_percent_delta(c.get("latency_cycles"), c.get("hardware_measured_cycles"))),
        ("roofline_compute_utilization(%)", lambda c: fmt_percent(c.get("roofline_compute_utilization"))),
        ("analytical_compute_utilization(%)", lambda c: fmt_percent(c.get("analytical_efficiency"))),
        ("measured_compute_utilization(%)", lambda c: fmt_percent(c.get("measured_efficiency"))),
        ("roofline_bandwidth_utilization(%)", lambda c: fmt_percent(c.get("roofline_bandwidth_utilization"))),
        ("analytical_bandwidth_utilization(%)", lambda c: fmt_percent(c.get("analytical_bandwidth_utilization"))),
        ("measured_bandwidth_utilization(%)", lambda c: fmt_percent(c.get("measured_bandwidth_utilization"))),
    ]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "ring_gemm 单算子性能表", "不同 size 对比；指标口径对齐 export_op_performance_summary.py / roofline compact CSV")
    add_box(slide, Inches(0.72), Inches(1.02), Inches(11.9), Inches(0.62),
            "对象：ring_gemm_fp16_fp16_fp16_0  |  mode=baseline  |  同地址映射思路下的三组 case 对比",
            fill=LIGHT_GREEN, line=GREEN, font_size=15, bold=True)

    x0 = Inches(0.82)
    y0 = Inches(1.82)
    c0 = Inches(4.25)
    c1 = Inches(2.4)
    c2 = Inches(2.4)
    c3 = Inches(2.4)
    row_h = Inches(0.37)

    add_box(slide, x0, y0, c0, row_h, "metric", fill=LIGHT_BLUE, line=BLUE, font_size=13, bold=True)
    add_box(slide, x0 + c0, y0, c1, row_h, cases[0]["label"], fill=LIGHT_BLUE, line=BLUE, font_size=12, bold=True)
    add_box(slide, x0 + c0 + c1, y0, c2, row_h, cases[1]["label"], fill=LIGHT_BLUE, line=BLUE, font_size=12, bold=True)
    add_box(slide, x0 + c0 + c1 + c2, y0, c3, row_h, cases[2]["label"], fill=LIGHT_BLUE, line=BLUE, font_size=11, bold=True)

    for idx, (metric, formatter) in enumerate(metrics, start=1):
        y = y0 + row_h * idx
        fill = WHITE if idx % 2 else LIGHT_GRAY
        add_box(slide, x0, y, c0, row_h, metric, fill=fill, line=GRAY, font_size=12, align=PP_ALIGN.LEFT)
        add_box(slide, x0 + c0, y, c1, row_h, formatter(cases[0]), fill=fill, line=GRAY, font_size=12)
        add_box(slide, x0 + c0 + c1, y, c2, row_h, formatter(cases[1]), fill=fill, line=GRAY, font_size=12)
        add_box(slide, x0 + c0 + c1 + c2, y, c3, row_h, formatter(cases[2]), fill=fill, line=GRAY, font_size=12)

    add_box(slide, Inches(0.88), Inches(5.78), Inches(2.9), Inches(0.9),
            "旧 size 结论\nmemory-bound，但已优于 measured\nerror = -4.76%",
            fill=LIGHT_RED, line=RED, font_size=14, bold=True)
    add_box(slide, Inches(3.98), Inches(5.78), Inches(3.15), Inches(0.9),
            f"新 size 结论\nlatency={int(current['latency_cycles'])}, measured={int(current['hardware_measured_cycles'])}\nerror = {fmt_percent_delta(current.get('latency_cycles'), current.get('hardware_measured_cycles'))}",
            fill=LIGHT_ORANGE, line=ORANGE, font_size=14, bold=True)
    add_box(slide, Inches(7.33), Inches(5.78), Inches(2.35), Inches(0.9),
            f"28-slice 结论\nlatency={int(case_28['latency_cycles'])}, measured={int(case_28['hardware_measured_cycles'])}\nerror = {fmt_percent_delta(case_28.get('latency_cycles'), case_28.get('hardware_measured_cycles'))}",
            fill=LIGHT_PURPLE, line=PURPLE, font_size=13, bold=True)
    add_box(slide, Inches(9.92), Inches(5.78), Inches(2.0), Inches(0.9),
            "对比观察\n28-slice global ring 下\ncompute 利用率更高，但最终仍是 memory-bound",
            fill=LIGHT_GREEN, line=GREEN, font_size=12, bold=True)
    add_footer(slide, "Slide 1 · 单算子性能对比表")


def main():
    output = Path("outputs/modeling/ring_gemm_performance_model.pptx")
    notes = Path("outputs/modeling/ring_gemm_performance_model_notes.md")
    build_deck(output, notes)
    print(output)


if __name__ == "__main__":
    main()
